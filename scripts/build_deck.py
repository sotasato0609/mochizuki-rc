#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
望月リソルゴルフクラブ 第2期会員募集｜運用分析レポート＆2026年8月配信提案

デザインフォーマット: 「営業本部共有_買取シミュレーター_2026-07.pptx」に準拠。
参照元から抽出した design token をそのまま使用する（.claude/skills/pptx/SKILL.md 参照）。
  - フォント: Meiryo 単一 ／ 角丸なし（すべて直角矩形）／ 影なし
  - ヘッダ文法: 縦バー+キッカー(11.5pt) → タイトル(27pt) → 短バー(334A62) + 全幅ヘアライン(E2E6EB)
  - フッタ: ページ番号のみ（罫線・社名なし）
  - カード: F2F4F7 + 上辺 0.05in のアクセントバー
  - 表: ヘッダ 405D7B/白 → 行は FFFFFF / EEF0F3 の交互、行高 0.58in

正確性プロトコル（.claude/rules/general.md）: 両論併記・証拠クラス明示・内部整合性。
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
import os

OUT_DIR = "output"
OUT = os.path.join(OUT_DIR, "望月リソルGC_運用分析レポート＆2026年8月配信提案.pptx")

# ══ 参照資料から抽出したカラーパレット ═══════════════════════
COLORS = {
    "primary":     RGBColor(0x40, 0x5D, 0x7B),
    "primary_dk":  RGBColor(0x33, 0x4A, 0x62),
    "accent":      RGBColor(0xA9, 0x69, 0x33),
    "gray":        RGBColor(0x6C, 0x75, 0x82),
    "ink":         RGBColor(0x2A, 0x2E, 0x33),
    "text":        RGBColor(0x45, 0x4B, 0x52),
    "muted":       RGBColor(0x8A, 0x92, 0x9C),
    "slate_lt":    RGBColor(0x8F, 0xA3, 0xB8),
    "surface":     RGBColor(0xF2, 0xF4, 0xF7),
    "surface_alt": RGBColor(0xEE, 0xF0, 0xF3),
    "border":      RGBColor(0xE2, 0xE6, 0xEB),
    "hairline":    RGBColor(0xD6, 0xDB, 0xE1),
    "bg":          RGBColor(0xFF, 0xFF, 0xFF),
}

# ══ タイポグラフィ（参照資料の実測値）═══════════════════════
FONT_JA = "Meiryo"
SZ = {
    "cover_title": 33, "cover_right": 22, "cover_sub": 15, "cover_kick": 13, "cover_meta": 11.5,
    "title": 27, "card_head": 15.5, "lead": 13.5, "body": 13, "body_s": 12,
    "kicker": 11.5, "cell": 11.5, "label": 10.5, "fine": 9.5, "fine_s": 8.5,
    "num": 27, "num_s": 22,
}

# ══ グリッド（参照資料の実測値）════════════════════════════
SW, SH = 13.333, 7.5
MX = 0.62
CW = 12.10
TITLE_X, TITLE_W = 0.58, 12.00
BODY_Y = 1.90
PAGE_Y = 7.12
G = 0.20                       # カード間ギャップ

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(SW), Inches(SH)
BLANK = prs.slide_layouts[6]
_st = {"page": 0, "layouts": []}
I = Inches


# ══ ヘルパ ══════════════════════════════════════════
def txt(slide, x, y, w, h, text, size, color=None, bold=False,
        align=PP_ALIGN.LEFT, line=1.35, anchor=MSO_ANCHOR.TOP, space_after=3):
    tb = slide.shapes.add_textbox(I(x), I(y), I(w), I(h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, ln in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln
        p.alignment = align
        p.line_spacing = line
        p.space_after = Pt(space_after)
        # 参考資料に合わせ、フォント指定は run 単位で行う
        for fnt in [p.font] + [r.font for r in p.runs]:
            fnt.name = FONT_JA
            fnt.size = Pt(size)
            fnt.bold = bold
            fnt.color.rgb = color or COLORS["text"]
    return tb


def rect(slide, x, y, w, h, fill):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, I(x), I(y), I(w), I(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    sp.line.fill.background()
    sp.shadow.inherit = False
    sp.text_frame.text = ""
    return sp


def card(slide, x, y, w, h, accent=None, fill=None):
    """F2F4F7 の面 + 上辺 0.05in のアクセントバー（参照資料の作法）"""
    rect(slide, x, y, w, h, fill or COLORS["surface"])
    if accent:
        rect(slide, x, y, w, 0.05, accent)


def header(slide, kicker, title, tag):
    """縦バー+キッカー → タイトル27pt → 短バー + 全幅ヘアライン"""
    _st["page"] += 1
    _st["layouts"].append(tag)
    L = _st["layouts"]
    if len(L) >= 3 and L[-1] == L[-2] == L[-3]:
        raise AssertionError(f"同一レイアウト {tag} が3枚連続（p{_st['page']}）")
    rect(slide, MX, 0.50, 0.09, 0.26, COLORS["primary"])
    txt(slide, 0.80, 0.48, 11.5, 0.30, kicker, SZ["kicker"], COLORS["primary"], True)
    txt(slide, TITLE_X, 0.80, TITLE_W, 0.62, title, SZ["title"], COLORS["ink"], True, line=1.15)
    rect(slide, MX, 1.53, 1.75, 0.05, COLORS["primary_dk"])
    rect(slide, MX, 1.55, CW, 0.01, COLORS["border"])
    txt(slide, 12.40, PAGE_Y, 0.70, 0.30, f"{_st['page']:02d}", SZ["label"],
        COLORS["muted"], align=PP_ALIGN.RIGHT)


def slide_new(kicker, title, tag="C", lead=None):
    s = prs.slides.add_slide(BLANK)
    header(s, kicker, title, tag)
    if lead:
        txt(s, MX, 1.68, CW, 0.3, lead, SZ["kicker"], COLORS["muted"])
    return s


def divider(no, title, sub):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, COLORS["primary"])
    rect(s, 0, 0, 0.06, SH, COLORS["accent"])
    txt(s, 1.25, 2.95, 10, 0.34, no, SZ["cover_kick"], COLORS["accent"], True)
    txt(s, 1.25, 3.42, 10.5, 0.7, title, SZ["cover_title"], COLORS["bg"], True, line=1.15)
    rect(s, 1.25, 4.42, 2.30, 0.04, COLORS["accent"])
    txt(s, 1.25, 4.72, 10.5, 0.5, sub, SZ["cover_sub"], COLORS["hairline"], line=1.4)
    _st["page"] += 1
    _st["layouts"].append("B")
    txt(s, 12.40, PAGE_Y, 0.70, 0.30, f"{_st['page']:02d}", SZ["label"],
        COLORS["slate_lt"], align=PP_ALIGN.RIGHT)
    return s


def table(slide, x, y, w, rows, col_w, row_h=0.44, head_h=0.46,
          fs=None, aligns=None, hl=None):
    """ヘッダ 405D7B/白 → 行 FFFFFF / EEF0F3 交互（参照資料の作法）"""
    fs = fs or SZ["cell"]
    hl = hl or []
    ncol = len(rows[0])
    tot = sum(col_w)
    col_w = [c / tot * w for c in col_w]
    aligns = aligns or [PP_ALIGN.LEFT] + [PP_ALIGN.RIGHT] * (ncol - 1)
    cy = y
    for ri, row in enumerate(rows):
        h = head_h if ri == 0 else row_h
        if ri == 0:
            bg = COLORS["primary"]
        elif ri in hl:
            bg = COLORS["surface"]
        else:
            bg = COLORS["bg"] if ri % 2 == 1 else COLORS["surface_alt"]
        rect(slide, x, cy, w, h, bg)
        cx = x
        for ci, cell in enumerate(row):
            v = str(cell)
            strong = v.startswith("**")
            v = v.replace("**", "")
            col = COLORS["bg"] if ri == 0 else (COLORS["primary"] if strong else COLORS["text"])
            txt(slide, cx + 0.12, cy, col_w[ci] - 0.24, h, v, fs, col,
                bold=(ri == 0) or strong,
                align=aligns[ci] if ci < len(aligns) else PP_ALIGN.LEFT,
                line=1.0, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
            cx += col_w[ci]
        cy += h
    return cy


def kpi(slide, x, y, w, h, value, label, note=None, accent=None, vcolor=None, vsize=None):
    card(slide, x, y, w, h, accent or COLORS["primary"])
    txt(slide, x + 0.22, y + 0.22, w - 0.44, 0.28, label, SZ["label"], COLORS["gray"], True)
    txt(slide, x + 0.22, y + 0.54, w - 0.44, 0.62, value, vsize or SZ["num"],
        vcolor or COLORS["primary"], True, line=1.1)
    if note:
        txt(slide, x + 0.22, y + h - 0.6, w - 0.44, 0.5, note, SZ["fine"], COLORS["muted"], line=1.3)


def panel(slide, x, y, w, h, label, title, body, accent=None, dark=False):
    """課題/方針の対比パネル（参照資料スライド3の作法）"""
    rect(slide, x, y, w, h, COLORS["primary"] if dark else COLORS["surface"])
    if accent and not dark:
        rect(slide, x, y, w, 0.05, accent)
    lc = COLORS["slate_lt"] if dark else COLORS["gray"]
    tc = COLORS["bg"] if dark else COLORS["ink"]
    bc = COLORS["bg"] if dark else COLORS["text"]
    txt(slide, x + 0.23, y + 0.22, w - 0.46, 0.3, label, SZ["label"], lc, True)
    txt(slide, x + 0.23, y + 0.48, w - 0.46, 0.4, title, 16, tc, True)
    txt(slide, x + 0.26, y + 1.02, w - 0.52, h - 1.24, body, SZ["body"], bc, line=1.5)


def bullets(items, mark="･"):
    return "\n".join(f"{mark}  {t}" for t in items)


def note_line(slide, y, s):
    txt(slide, MX, y, CW, 0.4, s, SZ["fine"], COLORS["muted"], line=1.35)


def chart(slide, kind, x, y, w, h, cats, series, colors, legend=False,
          labels=True, lblsize=11, lblcolor=None, gap=60, overlap=None):
    cd = CategoryChartData()
    cd.categories = cats
    for nm, vals in series:
        cd.add_series(nm, vals)
    gf = slide.shapes.add_chart(kind, I(x), I(y), I(w), I(h), cd)
    ch = gf.chart
    ch.font.size = Pt(10)
    ch.font.name = FONT_JA
    ch.font.color.rgb = COLORS["text"]
    if legend:
        ch.has_legend = True
        ch.legend.position = XL_LEGEND_POSITION.BOTTOM
        ch.legend.include_in_layout = False
        ch.legend.font.size = Pt(10)
    else:
        ch.has_legend = False
    try:
        va = ch.value_axis
        va.has_major_gridlines = True
        va.major_gridlines.format.line.color.rgb = COLORS["border"]
        va.major_gridlines.format.line.width = Pt(0.5)
        va.tick_labels.font.size = Pt(9.5)
        va.format.line.fill.background()
    except Exception:
        pass
    try:
        ca = ch.category_axis
        ca.tick_labels.font.size = Pt(9.5)
        ca.format.line.color.rgb = COLORS["border"]
    except Exception:
        pass
    pl = ch.plots[0]
    pl.gap_width = gap
    if overlap is not None:
        pl.overlap = overlap
    if labels:
        pl.has_data_labels = True
        pl.data_labels.font.size = Pt(lblsize)
        pl.data_labels.font.bold = True
        pl.data_labels.font.color.rgb = lblcolor or COLORS["ink"]
    for i, c in enumerate(colors):
        sr = pl.series[i]
        sr.format.fill.solid()
        sr.format.fill.fore_color.rgb = c
    return ch, pl


C3 = 3.90   # 3カラム幅
C2 = 5.95   # 2カラム幅
X3 = [MX, MX + C3 + G, MX + (C3 + G) * 2]
X2 = [MX, MX + C2 + G]

# ══════════════════════════════════════════════════
# 01  表紙
# ══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 8.75, 0, 4.58, SH, COLORS["primary"])
rect(s, 8.75, 0, 0.06, SH, COLORS["accent"])
txt(s, 0.90, 0.90, 5.0, 0.4, "望月リソルゴルフクラブ 第2期会員募集", SZ["cover_kick"],
    COLORS["primary"], True)
rect(s, 0.95, 1.55, 6.90, 0.02, COLORS["hairline"])
txt(s, 0.90, 2.35, 7.4, 1.7, "WEB広告\n運用分析レポート", SZ["cover_title"],
    COLORS["ink"], True, line=1.2)
txt(s, 9.15, 2.90, 3.6, 1.8,
    "Mochizuki Resol\nGolf Club\n\n運用分析 ＆\n8月配信提案",
    SZ["cover_right"], COLORS["bg"], True, line=1.35)
txt(s, 0.92, 4.20, 7.0, 0.6, "および 2026年8月 配信提案", SZ["cover_sub"], COLORS["text"])
rect(s, 0.95, 5.95, 2.30, 0.04, COLORS["primary_dk"])
txt(s, 0.92, 6.20, 7.0, 0.9, "2026年7月 ／ ゲンダイエージェンシー株式会社",
    SZ["cover_meta"], COLORS["gray"])
_st["page"] = 1
_st["layouts"].append("A")

# ══ 02  本レポートについて ═══════════════════════════
s = slide_new("はじめに", "本レポートについて", tag="C")
txt(s, MX, BODY_Y, CW, 0.3,
    "成果が伸びていない要因を特定するため、第1期にさかのぼって全データを検証しました。",
    SZ["body"], COLORS["text"])
rows = [
    ["調査対象", "期間", "内容"],
    ["お問い合わせ実データ", "2025年8月〜2026年7月", "通知メールを1件ずつ精査。スパム・テスト送信を除外"],
    ["Google Analytics 4", "2024年10月〜2026年7月", "セッション・地域・デバイス・流入元・遷移先"],
    ["Google広告", "2025年6月〜2026年7月", "キャンペーン別実績、検索テーマ、エリア設定"],
    ["Meta広告", "2024年10月〜2026年7月", "月別のエリア設定の変遷、オーディエンス別成績"],
    ["DV360", "2024年11月〜2026年7月", "全キャンペーン実績、配信面別データ 9,833行"],
    ["LINEヤフー広告", "2024年10月〜2025年6月", "配信実績（表示・クリック・クリック率）"],
    ["ご請求データ", "2024年9月〜2026年7月", "媒体別・月別の費用構成"],
]
table(s, MX, BODY_Y + 0.45, CW, rows, [2.7, 2.6, 6.8],
      aligns=[PP_ALIGN.LEFT] * 3, row_h=0.42, head_h=0.44)
card(s, MX, 6.16, CW, 0.72, COLORS["accent"])
txt(s, MX + 0.24, 6.16, CW - 0.48, 0.72,
    "推測で数字を埋めていません。確認できなかった項目は「未確認」と明記し、件数が少ない項目は断定せず「方向性」として扱っています。",
    SZ["body"], COLORS["text"], anchor=MSO_ANCHOR.MIDDLE)

# ══ 03  サマリー（カード6） ═══════════════════════════
s = slide_new("概要", "分析で分かったこと", tag="E")
txt(s, MX, BODY_Y, CW, 0.3,
    "ご予算を30%増やし集客を1.8倍にしたにもかかわらず、お問い合わせは1/3以下に減少しました。",
    SZ["body"], COLORS["text"])
F = [("01", "成果の53%は、ゴルフ場から\n概ね70km圏内の方でした", COLORS["primary"]),
     ("02", "配信エリアと成果は、月単位で\nきれいに対応していました", COLORS["accent"]),
     ("03", "成果が出ていた時期は、\n3媒体がバランスしていました", COLORS["gray"]),
     ("04", "媒体によってクリックがサイトに\n届く割合が2.7倍違います", COLORS["primary"]),
     ("05", "検索クリックの70%が、すでに\n貴クラブをご存知の方でした", COLORS["accent"]),
     ("06", "じっくり読まれているのはPC。\nしかし流入の96%はスマホです", COLORS["gray"])]
for i, (no, body, acc) in enumerate(F):
    x = X3[i % 3]
    y = BODY_Y + 0.5 + (i // 3) * (1.62 + G)
    card(s, x, y, C3, 1.62, acc)
    txt(s, x + 0.22, y + 0.24, C3 - 0.44, 0.28, no, SZ["label"], acc, True)
    txt(s, x + 0.22, y + 0.6, C3 - 0.44, 0.8, body, SZ["card_head"], COLORS["ink"], True, line=1.4)
note_line(s, 6.62,
          "証拠クラス〈実測〉：全項目が各媒体API・GA4・お問い合わせ実データからの実測値です。業界一般値・当社仮説は含みません。")

# ══ 04  Divider ═════════════════════════════════
divider("PART 1", "成果の全体像",
        "第1期（2025年8月〜2026年3月）と第2期（2026年5月〜7月）を、同じ基準で比較します。")

# ══ 05  KPI ═════════════════════════════════════
s = slide_new("1 成果", "第1期・第2期の比較", tag="F")
K = [("+30.0%", "月あたりご請求額", "¥251,250 → ¥326,667", COLORS["primary"], COLORS["primary"]),
     ("1.8倍", "月あたりサイト集客", "6,427 → 11,715", COLORS["primary"], COLORS["primary"]),
     ("−89%", "お問い合わせ", "19件 → 2件", COLORS["accent"], COLORS["accent"]),
     ("4.63倍", "獲得単価", "¥105,789 → ¥490,000", COLORS["accent"], COLORS["accent"])]
w4 = (CW - G * 3) / 4
for i, (v, l, n, ac, vc) in enumerate(K):
    kpi(s, MX + i * (w4 + G), BODY_Y, w4, 1.9, v, l, n, accent=ac, vcolor=vc)
rows = [
    ["", "第1期  2025年8月〜2026年3月（8ヶ月）", "第2期  2026年5月〜7月（3ヶ月）"],
    ["ご請求額 合計", "¥2,010,000", "¥980,000"],
    ["サイトセッション", "51,414", "31,630"],
    ["**お問い合わせ", "**19件", "**2件"],
    ["お問い合わせ率（CVR）", "0.037%", "0.006%"],
    ["**獲得単価", "**¥105,789", "**¥490,000"],
]
table(s, MX, BODY_Y + 2.12, CW, rows, [4.2, 4.0, 3.9], hl=[3, 5], row_h=0.42, head_h=0.44)
note_line(s, 6.5,
          "〈実測〉ご請求額＝弊社基幹システム／集客＝GA4／お問い合わせ＝通知メールの実件数（スパム・テスト送信を除外）。\n"
          "比較期間はお問い合わせの実データが確認できている期間に揃え、2026年4月（LP刷新期間）は除外しています。")

# ══ 06  月次推移 ═════════════════════════════════
s = slide_new("1 成果", "月次推移 ── 最良月は2026年3月", tag="G",
              lead="第1期は月を追うごとに改善していました。2026年1月2件 → 2月3件 → 3月5件。")
ch, pl = chart(s, XL_CHART_TYPE.COLUMN_CLUSTERED, MX, BODY_Y + 0.24, CW, 2.95,
               ["25/8", "25/9", "25/10", "25/11", "25/12", "26/1", "26/2", "26/3", "26/5", "26/6", "26/7"],
               [("お問い合わせ件数", (1, 3, 3, 0, 2, 2, 3, 5, 1, 1, 0))],
               [COLORS["primary"]], gap=55)
pl.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
p_ = pl.series[0].points[7]
p_.format.fill.solid()
p_.format.fill.fore_color.rgb = COLORS["accent"]
rows = [
    ["月", "25/8", "25/9", "25/10", "25/11", "25/12", "26/1", "26/2", "26/3", "26/5", "26/6", "26/7"],
    ["ご請求額（千円）", "300", "300", "280", "100", "150", "330", "300", "**250", "300", "380", "300"],
    ["獲得単価（千円）", "300", "100", "93", "—", "75", "165", "100", "**50", "300", "380", "—"],
]
table(s, MX, BODY_Y + 3.4, CW, rows, [2.2] + [0.9] * 11,
      aligns=[PP_ALIGN.LEFT] + [PP_ALIGN.RIGHT] * 11, row_h=0.4, head_h=0.4, fs=10.5)
note_line(s, 6.62,
          "〈実測〉オレンジは全期間で最良の月。〈両論併記〉2025年11月の0件は予算を¥100,000に縮小した月である点を付記します。")

# ══ 07  見込み度 ═════════════════════════════════
s = slide_new("1 成果", "件数だけでなく、見込み度も下がっています", tag="D")
kpi(s, X2[0], BODY_Y, C2, 1.85, "68%", "第1期（19件）｜具体的に検討している",
    "13件が「具体的に検討している」と回答", accent=COLORS["primary"], vsize=36)
kpi(s, X2[1], BODY_Y, C2, 1.85, "0%", "第2期（2件）｜具体的に検討している",
    "「興味がある」1件、来場予定の方からのご連絡1件", accent=COLORS["accent"],
    vcolor=COLORS["accent"], vsize=36)
rows = [
    ["会員権について", "第1期（19件）", "第2期（2件）"],
    ["**具体的に検討している", "**13件（68%）", "**0件"],
    ["興味がある", "7件", "1件"],
    ["未選択・その他", "—", "1件"],
]
table(s, MX, BODY_Y + 2.1, CW, rows, [5.6, 3.3, 3.2], hl=[1], row_h=0.46, head_h=0.46)
note_line(s, 6.14,
          "〈実測〉お問い合わせ通知メールの回答項目を1件ずつ分類。母集団は第1期19件・第2期2件（いずれも全件）。\n"
          "〈方向性〉第2期は母数が2件と少なく、比率は断定ではなく方向性としてご理解ください。")

# ══ 08  Divider ═════════════════════════════════
divider("PART 2", "誰が申し込んでいるのか",
        "第1期のお問い合わせ19件を、1件ずつ内容を確認して分類しました。")

# ══ 09  入口 ════════════════════════════════════
s = slide_new("2 顧客", "入口は「資料請求」です", tag="F")
kpi(s, MX, BODY_Y, 4.0, 2.3, "95%", "お問い合わせ種別｜資料のご請求",
    "19件中18件", accent=COLORS["accent"], vcolor=COLORS["accent"], vsize=44)
rows = [
    ["お問い合わせ種別", "件数", "比率"],
    ["**資料のご請求", "**18件", "**95%"],
    ["視察プレーのお問い合わせ", "3件", "16%"],
    ["その他", "1件", "5%"],
]
table(s, MX + 4.0 + G, BODY_Y, CW - 4.0 - G, rows, [4.4, 1.6, 1.6], hl=[1], row_h=0.5, head_h=0.5)
txt(s, MX + 4.0 + G, BODY_Y + 2.06, CW - 4.0 - G, 0.3,
    "※ 複数選択のため合計は100%を超えます", SZ["fine"], COLORS["muted"])
panel(s, MX, BODY_Y + 2.62, CW, 1.42, "示唆", "入口は資料請求に置くべきです",
      "視察プレーの直接お申し込みは3件のみ。大半の方は、まず資料を請求し、そこから検討に入られています。\n"
      "8月に視察プレーを訴求する場合も、いきなり来場を求めるのではなく、資料請求の受け皿を用意しつつ期限を訴えるのが自然です。",
      accent=COLORS["primary"])
note_line(s, 6.24,
          "〈実測〉母集団＝第1期のお問い合わせ全19件。〈当社仮説〉示唆部分は実測の解釈であり、検証はこれからです。")

# ══ 10  距離帯 ══════════════════════════════════
s = slide_new("2 顧客", "成果の53%は、ゴルフ場から概ね70km圏内でした", tag="D")
ch, pl = chart(s, XL_CHART_TYPE.BAR_CLUSTERED, MX, BODY_Y, 7.3, 3.4,
               ["概ね30km圏\n佐久・小諸・軽井沢・御代田", "概ね40km圏\n上田市",
                "概ね70km圏\n群馬県高崎市", "概ね90km圏\n長野県松本市",
                "首都圏\n東京・神奈川・千葉"],
               [("件数", (7, 2, 1, 1, 8))], [COLORS["primary"]], gap=45, lblsize=12)
for idx in (0, 1, 2):
    p_ = pl.series[0].points[idx]
    p_.format.fill.solid()
    p_.format.fill.fore_color.rgb = COLORS["accent"]
rx, rw = MX + 7.3 + G, CW - 7.3 - G
kpi(s, rx, BODY_Y, rw, 1.62, "53%", "概ね70km圏内の累計",
    "19件中10件。佐久市内2件・隣接町村5件", accent=COLORS["accent"],
    vcolor=COLORS["accent"], vsize=36)
rows = [
    ["距離帯", "件数", "累計"],
    ["**概ね30km圏", "**7件", "**37%"],
    ["概ね40km圏", "2件", "47%"],
    ["**概ね70km圏", "**1件", "**53%"],
    ["概ね90km圏", "1件", "58%"],
    ["首都圏", "8件", "100%"],
]
table(s, rx, BODY_Y + 1.82, rw, rows, [2.1, 0.9, 0.9], row_h=0.3, head_h=0.32, fs=10.5)
note_line(s, 6.0,
          "〈実測〉距離帯はお問い合わせ記載の郵便番号・住所から分類。\n"
          "〈両論併記〉首都圏も8件（42%）あり、次ページの通り検討度はむしろ高い水準です。")

# ══ 11  地元と首都圏 ════════════════════════════
s = slide_new("2 顧客", "地元は件数、首都圏は検討度", tag="D",
              lead="地元と首都圏では役割が違います。一律に扱うべきではありません。")
panel(s, X2[0], BODY_Y + 0.32, C2, 2.5, "地元（長野・群馬）",
      "11件 ／ うち具体的に検討 64%",
      "件数が多い。ゴルフ場のすぐ近くから着実にお問い合わせが来ています。\n"
      "佐久市内が2件、隣接する町村が5件。",
      accent=COLORS["primary"])
panel(s, X2[1], BODY_Y + 0.32, C2, 2.5, "首都圏（東京・神奈川・千葉）",
      "8件 ／ うち具体的に検討 75%",
      "件数は少ないが、検討度が高い。「遠い」ことを自覚したうえでご検討されています。",
      dark=True)
panel(s, MX, BODY_Y + 3.06, CW, 1.2, "結論", "別枠に分けて成果を測るべきです",
      "唯一の遠方からの資料請求（神奈川県茅ヶ崎市）は「自宅から遠方のため、まずは資料をもとに検討を始めたい」と述べられていました。",
      accent=COLORS["accent"])
note_line(s, 6.44,
          "〈実測〉母集団＝19件（地元11件・首都圏8件）。〈方向性〉母数が少ないため、比率の差は断定ではなく方向性です。")

# ══ 12  Divider ═════════════════════════════════
divider("PART 3", "配信の分析",
        "エリア設定・媒体構成が、成果とどう対応していたかを月単位で検証しました。")

# ══ 13  ミスマッチ ══════════════════════════════
s = slide_new("3 配信", "成果を生む地域に、集客が向いていません", tag="D")
chart(s, XL_CHART_TYPE.COLUMN_CLUSTERED, MX, BODY_Y, 7.3, 3.5,
      ["長野県", "群馬県", "東京都", "神奈川県"],
      [("第1期の成果構成比", (53, 5, 32, 5)),
       ("2026年7月の集客構成比", (16.7, 2.6, 35.7, 13.6))],
      [COLORS["primary"], COLORS["border"]], legend=True, gap=70, lblsize=10)
rx, rw = MX + 7.3 + G, CW - 7.3 - G
rows = [
    ["地域", "成果", "集客", "差"],
    ["**長野県", "**53%", "**16.7%", "**−36pt"],
    ["群馬県", "5%", "2.6%", "−2pt"],
    ["東京都", "32%", "35.7%", "+4pt"],
    ["神奈川県", "5%", "13.6%", "+9pt"],
]
table(s, rx, BODY_Y, rw, rows, [1.5, 0.85, 0.9, 0.95], hl=[1], row_h=0.42, head_h=0.44)
panel(s, rx, BODY_Y + 2.28, rw, 1.22, "所見", "向き先がズレています",
      "成果の半分以上を生む長野県に、集客の2割弱しか向けていません。",
      accent=COLORS["accent"])
note_line(s, 6.0,
          "〈実測〉成果構成比＝第1期19件の居住地分布／集客構成比＝GA4の2026年7月セッション。\n"
          "分母が異なるため同一指標の比較ではなく、「向き先のズレ」を示すものです。")

# ══ 14  地域別の質 ══════════════════════════════
s = slide_new("3 配信", "長野からの流入は、最も質が高い", tag="G",
              lead="第2期（2026年5月〜7月）の地域別。集客量が最も少ない長野が、最も読まれています。")
rows = [
    ["地域", "セッション", "構成比", "エンゲージ率", "平均滞在"],
    ["東京都", "10,262", "35.7%", "21.1%", "3.2秒"],
    ["**長野県", "**4,802", "**16.7%", "**25.2%", "**5.8秒"],
    ["神奈川県", "3,895", "13.6%", "22.9%", "2.8秒"],
    ["埼玉県", "1,918", "6.7%", "22.7%", "2.5秒"],
    ["大阪府", "1,360", "4.7%", "24.0%", "2.5秒"],
]
table(s, MX, BODY_Y + 0.28, CW, rows, [3.3, 2.2, 2.0, 2.2, 2.0], hl=[2], row_h=0.5, head_h=0.5)
panel(s, MX, BODY_Y + 3.42, CW, 1.16, "所見", "量は最も少なく、質は最も高い",
      "長野県は集客構成比では東京の半分以下。しかし平均滞在は東京の1.8倍、エンゲージ率も最も高い水準です。",
      accent=COLORS["primary"])
note_line(s, 6.62,
          "〈実測〉GA4（2026年5月1日〜7月21日）。平均滞在は総エンゲージメント時間÷セッション数で算出。")

# ══ 15  エリア設定と成果 ═════════════════════════
s = slide_new("3 配信", "配信エリアの設定と成果は、月単位で対応していました", tag="D")
rows = [
    ["月", "Meta広告の配信エリア", "お問い合わせ"],
    ["2025年9月", "東京23区", "3件"],
    ["2025年10月", "ゴルフ場から半径40km", "3件"],
    ["2025年11月", "Meta配信なし", "0件"],
    ["2025年12月", "東京23区", "2件"],
    ["2026年1月", "半径24km", "2件"],
    ["2026年2月", "半径48〜49km", "3件"],
    ["**2026年3月", "**半径48〜49km", "**5件（最良）"],
    ["2026年6月", "半径50km", "1件"],
    ["**2026年7月", "**東京・神奈川・長野の3都県に拡大", "**0件"],
]
table(s, MX, BODY_Y, 7.3, rows, [1.9, 3.7, 1.7],
      aligns=[PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.RIGHT], hl=[7], row_h=0.4, head_h=0.42)
rx, rw = MX + 7.3 + G, CW - 7.3 - G
kpi(s, rx, BODY_Y, rw, 1.4, "12件", "半径24〜50km圏で運用していた月の合計",
    accent=COLORS["primary"], vsize=30)
kpi(s, rx, BODY_Y + 1.6, rw, 1.4, "0件", "3都県に広げた2026年7月",
    accent=COLORS["accent"], vcolor=COLORS["accent"], vsize=30)
panel(s, rx, BODY_Y + 3.2, rw, 1.34, "留保", "因果は証明できません",
      "7月1日にはエリア拡大に加え、Google広告のP-MAX移行・DV360停止も同時に発生しています。",
      accent=COLORS["gray"])
note_line(s, 6.62,
          "〈実測〉Meta Marketing API から月別の広告セット設定を取得。〈両論併記〉因果の留保を右下に明記しています。")

# ══ 16  媒体構成 ════════════════════════════════
s = slide_new("3 配信", "成果が出ていた時期は、3媒体がバランスしていました", tag="G")
chart(s, XL_CHART_TYPE.COLUMN_STACKED, MX, BODY_Y, 7.6, 3.6,
      ["25年9月\n3件", "25年10月\n3件", "26年2月\n3件", "26年3月\n5件",
       "26年5月\n1件", "26年6月\n1件", "26年7月\n0件"],
      [("Google広告", (33, 40, 50, 40, 17, 14, 17)),
       ("Meta広告", (33, 40, 33, 20, 33, 14, 83)),
       ("DV360", (33, 20, 17, 40, 50, 71, 0))],
      [COLORS["primary"], COLORS["accent"], COLORS["gray"]],
      legend=True, gap=55, lblsize=9, lblcolor=COLORS["bg"])
rx, rw = MX + 7.6 + G, CW - 7.6 - G
rows = [
    ["", "Google", "Meta", "DV360", "件数/月"],
    ["**第1期 月平均", "**38%", "33%", "28%", "**2.4件"],
    ["**第2期 月平均", "**16%", "42%", "42%", "**0.67件"],
]
table(s, rx, BODY_Y, rw, rows, [1.5, 0.75, 0.7, 0.8, 0.9], row_h=0.5, head_h=0.42, fs=10.5)
panel(s, rx, BODY_Y + 1.5, rw, 2.1, "所見", "Google広告が半減しました",
      "月¥75,000から¥40,000へ約半減。\n2026年7月はMeta広告に83%が集中し、0件でした。\n"
      "最良月でさえMeta広告は¥40,000です。",
      accent=COLORS["accent"])
note_line(s, 6.0,
          "〈実測〉構成比＝ご請求データの配信費ベース。金額の多寡ではなく、媒体の組み合わせが成果に対応していました。")

# ══ 17  Divider ═════════════════════════════════
divider("PART 4", "配信品質の分析",
        "買ったクリックは、本当にサイトに届いているか。新しい評価軸で媒体を比較します。")

# ══ 18  サイト到達率 ════════════════════════════
s = slide_new("4 品質", "新しい評価軸 ── サイト到達率", tag="F")
card(s, MX, BODY_Y, CW, 0.72, COLORS["primary"])
txt(s, MX, BODY_Y, CW, 0.72, "サイト到達率　＝　GA4のセッション数　÷　広告のクリック数",
    16, COLORS["ink"], True, PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
txt(s, MX, BODY_Y + 0.92, CW, 0.4,
    "広告レポートの「クリック数」は、そのままサイト訪問数にはなりません。誤タップや、読み込み前の離脱があるためです。",
    SZ["body"], COLORS["text"])
kpi(s, X2[0], BODY_Y + 1.44, C2, 1.9, "43.0%",
    "LINEヤフー広告（2024年10月〜2025年5月）",
    "クリック率 0.12%　／　健全な水準", accent=COLORS["primary"], vsize=38)
kpi(s, X2[1], BODY_Y + 1.44, C2, 1.9, "15.9%", "DV360（2026年5〜6月）",
    "クリック率 11.30%　／　誤タップ主体", accent=COLORS["accent"],
    vcolor=COLORS["accent"], vsize=38)
note_line(s, 6.0,
          "〈実測〉クリック数＝各媒体API／セッション数＝GA4。〈業界一般〉ディスプレイ広告のクリック率は通常0.1〜0.5%程度とされます。\n"
          "〈両論併記〉DV360にも良質な配信面は存在します（次ページ）。")

# ══ 19  DV360 配信面 ════════════════════════════
s = slide_new("4 品質", "DV360の配信面を1面ずつ分析しました", tag="G",
              lead="2026年5〜6月の全配信面 9,833行を取得し、GA4と突き合わせました。")
rows = [
    ["配信面", "表示", "クリック", "クリック率", "サイト到達率", "平均滞在"],
    ["trilltrill.jp", "88,947", "7,574", "8.5%", "**18.5%", "**0.0秒"],
    ["SmartNews（iOSアプリ）", "23,728", "4,261", "**18.0%", "77.6%", "**0.9秒"],
    ["nlab.itmedia.co.jp", "25,312", "1,564", "6.2%", "12.6%", "0.8秒"],
    ["mamastar.jp", "23,619", "1,206", "5.1%", "**4.1%", "0.3秒"],
    ["MILE mobile（Androidアプリ）", "—", "—", "—", "—", "**28.7秒"],
    ["**全体", "**672,105", "**43,723", "**6.51%", "**15.9%", "—"],
]
table(s, MX, BODY_Y + 0.24, CW, rows, [3.3, 1.7, 1.6, 1.6, 1.9, 1.5],
      hl=[5, 6], row_h=0.4, head_h=0.42)
panel(s, MX, BODY_Y + 3.16, CW, 1.2, "所見", "面ごとの選別が有効です",
      "クリック率5%超の配信面が65面あり、クリックの61%・配信費の47%を占めています。"
      "一方でMILE mobileは平均滞在28.7秒と良質で、アプリ面が一律に悪いわけではありません。",
      accent=COLORS["primary"])
note_line(s, 6.5,
          "〈実測〉DV360 Bid Manager API の配信面別レポート（9,833行）をGA4と突合。\n"
          "※アプリ内ブラウザでは滞在時間が過少計測されうるため、判断の主軸は到達率に置いています。")

# ══ 20  LINEヤフー ══════════════════════════════
s = slide_new("4 品質", "LINEヤフー広告は、最も質の高い到達を実現していました", tag="D")
rows = [
    ["配信月", "表示回数", "クリック", "クリック率"],
    ["2024年10月", "1,657,045", "2,944", "0.18%"],
    ["2024年11月", "1,397,408", "2,544", "0.18%"],
    ["2025年5月", "3,738,868", "2,772", "0.07%"],
    ["**合計", "**6,793,326", "**8,260", "**0.12%"],
]
table(s, MX, BODY_Y, 7.0, rows, [2.1, 2.0, 1.5, 1.5], hl=[4], row_h=0.48, head_h=0.48)
rx, rw = MX + 7.0 + G, CW - 7.0 - G
kpi(s, rx, BODY_Y, rw, 1.5, "43.0%", "サイト到達率", "DV360（15.9%）の2.7倍",
    accent=COLORS["primary"], vsize=36)
kpi(s, rx, BODY_Y + 1.7, rw, 1.32, "680万回", "累計表示回数", accent=COLORS["primary"], vsize=28)
panel(s, MX, BODY_Y + 2.6, 7.0, 1.36, "再開について", "開設手続きなしで8月1日から配信できます",
      "貴クラブ専用アカウントが稼働可能な状態で残っています。",
      accent=COLORS["accent"])
note_line(s, 6.2,
          "〈実測〉LINEヤフー広告API。〈両論併記〉この期間のコンバージョンは0件でした。ただし当時コンバージョン計測が\n"
          "設定されていたかは未確認のため、成果の有無は断定できません。再開時は計測を設定したうえで配信します。")

# ══ 21  検索の分析 ══════════════════════════════
s = slide_new("4 品質", "検索クリックの70%が、すでに貴クラブをご存知の方でした", tag="G")
rows = [
    ["検索テーマ（2026年7月・Google広告）", "クリック", "比率"],
    ["**望月 リソル ゴルフ クラブ", "**35", "**44%"],
    ["**リソル 望月", "**11", "**14%"],
    ["**望月 東急（旧名称）", "**9", "**11%"],
    ["その他", "24", "30%"],
]
table(s, MX, BODY_Y, 6.9, rows, [4.0, 1.5, 1.4], hl=[1, 2, 3], row_h=0.46, head_h=0.46)
rx = MX + 6.9 + G
hw = (CW - 6.9 - G * 2) / 2
kpi(s, rx, BODY_Y, hw, 2.36, "70%", "クラブ名での検索", "検索クリックに占める比率",
    accent=COLORS["accent"], vcolor=COLORS["accent"], vsize=38)
kpi(s, rx + hw + G, BODY_Y, hw, 2.36, "42", "「ゴルフ 会員権」へのクリック（第1期）",
    "本命の検索にはほとんど出せていません", accent=COLORS["gray"],
    vcolor=COLORS["gray"], vsize=38)
panel(s, MX, BODY_Y + 2.56, CW, 1.4, "所見", "新規開拓に振り向ける余地があります",
      "第1期の検索広告では近隣ゴルフ場の名称にも出稿していましたが、そこからのお問い合わせは0件でした。\n"
      "クラブ名での検索は、広告を出さなくても自然検索で到達できます。",
      accent=COLORS["primary"])
note_line(s, 6.16,
          "〈実測〉Google Ads API の検索テーマレポート。第1期は検索語句ベースで費用の55.8%がクラブ名検索でした。")

# ══ 22  サイト内行動 ════════════════════════════
s = slide_new("4 品質", "サイト内行動の分析", tag="C")
txt(s, MX, BODY_Y, CW, 0.3,
    "① じっくり読まれているのはパソコン。しかし流入の96%はスマートフォンです",
    14, COLORS["primary"], True)
rows = [
    ["デバイス", "セッション", "構成比", "エンゲージ率", "平均滞在"],
    ["スマートフォン", "30,377", "**95.9%", "22.4%", "**2.5秒"],
    ["**パソコン", "**964", "**3.0%", "**47.5%", "**28.5秒"],
    ["タブレット", "326", "1.0%", "18.7%", "5.5秒"],
]
table(s, MX, BODY_Y + 0.36, 7.2, rows, [2.0, 1.5, 1.2, 1.5, 1.3],
      hl=[2], row_h=0.32, head_h=0.32, fs=10.5)
panel(s, MX + 7.2 + G, BODY_Y + 0.36, CW - 7.2 - G, 1.28, "所見",
      "スマートフォンの流入品質が低い",
      "パソコンは流入の3%ですが、平均滞在はスマートフォンの11.4倍です。",
      accent=COLORS["accent"])
txt(s, MX, BODY_Y + 1.82, CW, 0.3, "② 会員募集の導線が二重になっています",
    14, COLORS["primary"], True)
rows2 = [
    ["期間", "LPから公式サイトへの遷移", "お問い合わせ"],
    ["第1期（8ヶ月）", "767件", "19件"],
    ["第2期（3ヶ月）", "84件", "2件"],
    ["**全期間累計", "**1,430件", "—"],
]
table(s, MX, BODY_Y + 2.18, 5.6, rows2, [2.2, 2.4, 1.5], hl=[3],
      row_h=0.32, head_h=0.32, fs=10.5)
txt(s, MX + 5.6 + G, BODY_Y + 2.18, CW - 5.6 - G, 1.28,
    bullets(["LPと公式サイトで同じ会員募集を掲載しており、リンクは一方通行です",
             "ただし公式サイトにもお問い合わせフォームがあり、そちらで成果が発生している可能性があります",
             "「損失」とは決めつけられません"]),
    SZ["body_s"], COLORS["text"], line=1.5)
txt(s, MX, BODY_Y + 3.64, CW, 0.3, "③ 検討中の方への再接触が不足しています",
    14, COLORS["primary"], True)
rows3 = [
    ["Meta広告（直近30日）", "予算構成比", "クリック率", "LP到達率", "接触回数"],
    ["興味関心（新規開拓）", "**91.0%", "1.84%", "82.3%", "1.87回"],
    ["**リターゲティング（検討層）", "**9.0%", "**2.89%", "**85.7%", "2.12回"],
]
table(s, MX, BODY_Y + 4.0, CW, rows3, [3.5, 2.1, 2.0, 2.0, 2.1],
      hl=[2], row_h=0.32, head_h=0.32, fs=10.5)

# ══ 23  Divider ═════════════════════════════════
divider("PART 5", "2026年8月 配信提案",
        "分析から導かれた3つの課題と、視察プレー最終月に向けた打ち手をご提案します。")

# ══ 24  課題と打ち手 ════════════════════════════
s = slide_new("5 提案", "3つの課題と打ち手", tag="E",
              lead="いずれも「広告を増やす／減らす」ではなく、「向き先を変える」ことで対応できます。")
items = [
    ("課題 01", "配信エリアが広すぎる", "成果の53%が概ね70km圏／長野の集客は16.7%",
     "配信エリアを半径50〜80km圏へ戻す", "東京は停止せず別枠に分離し、エリア別に成果を測る"),
    ("課題 02", "検討層への再接触が足りていない", "クリック率は1.57倍なのに予算9%・接触2.12回",
     "リターゲティングを強化する", "比率を9%→20〜25%、接触回数4〜5回／オーディエンス拡張"),
    ("課題 03", "配信面と媒体構成に偏りがある", "DV360の到達率15.9%／7月はMeta83%の一極集中",
     "配信面を整理し、媒体構成を戻す", "65面を除外して再開／Google広告を第1期水準へ／LINEヤフー再開"),
]
y = BODY_Y + 0.34
lw = 5.35
for no, ttl, ev, act, det in items:
    card(s, MX, y, lw, 1.34, COLORS["accent"])
    txt(s, MX + 0.22, y + 0.22, 1.3, 0.24, no, SZ["label"], COLORS["accent"], True)
    txt(s, MX + 0.22, y + 0.5, lw - 0.44, 0.3, ttl, 14.5, COLORS["ink"], True)
    txt(s, MX + 0.22, y + 0.86, lw - 0.44, 0.34, ev, SZ["fine"], COLORS["muted"])
    txt(s, MX + lw + 0.06, y + 0.46, 0.42, 0.34, "▶", 14, COLORS["gray"], True, PP_ALIGN.CENTER)
    ax = MX + lw + 0.55
    rect(s, ax, y, CW - lw - 0.55, 1.34, COLORS["primary"])
    txt(s, ax + 0.24, y + 0.28, CW - lw - 1.03, 0.3, act, 14.5, COLORS["bg"], True)
    txt(s, ax + 0.24, y + 0.7, CW - lw - 1.03, 0.5, det, SZ["fine"], COLORS["hairline"], line=1.4)
    y += 1.52

# ══ 25  配信プラン ══════════════════════════════
s = slide_new("5 提案", "2026年8月 配信プラン", tag="G",
              lead="視察プレー期間（7月11日〜8月30日）の最終月。期間限定オファーの刈り取り月として設計します。")
rows = [
    ["媒体", "役割・設定", "配信費", "構成比"],
    ["Meta広告", "半径50〜80km（55%）／東京別枠（20%）／リターゲティング（25%）", "¥200,000", "47.6%"],
    ["Google広告（P-MAX）", "半径50〜80km＋東京別枠／クラブ名検索を除外", "¥80,000", "19.0%"],
    ["DV360", "競合店GEO／クリック率5%超の65面を除外", "¥80,000", "19.0%"],
    ["LINEヤフー広告", "長野県東信＋群馬県西部／既存アカウントを再開", "¥60,000", "14.3%"],
    ["配信費 小計", "", "¥420,000", ""],
    ["運用代行費", "", "¥105,000", ""],
    ["**合計", "", "**¥525,000", "**対7月 +75%"],
]
table(s, MX, BODY_Y + 0.24, CW, rows, [2.6, 5.9, 1.8, 1.6],
      aligns=[PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT],
      hl=[7], row_h=0.4, head_h=0.42)
rows2 = [
    ["", "Meta", "Google", "DV360", "LINEヤフー"],
    ["2026年7月（実績）", "**83.3%", "16.7%", "0%", "—"],
    ["**2026年8月（提案）", "**47.6%", "19.0%", "19.0%", "14.3%"],
]
table(s, MX, BODY_Y + 3.62, 6.5, rows2, [2.1, 1.1, 1.15, 1.05, 1.3],
      hl=[2], row_h=0.38, head_h=0.36, fs=10.5)
panel(s, MX + 6.5 + G, BODY_Y + 3.62, CW - 6.5 - G, 1.12, "補足",
      "Meta広告は減額しません",
      "内訳を組み替えて一極集中（83.3%）を47.6%まで下げます。",
      accent=COLORS["primary"])
note_line(s, 6.66, "〈実測〉7月実績＝ご請求データ。〈提案〉8月の構成比は本プランに基づく計画値です。")

# ══ 26  KPI ═════════════════════════════════════
s = slide_new("5 提案", "8月のKPI", tag="F")
k = [("3〜5件", "お問い合わせ", "第1期の月平均2.4件／最良月5件", COLORS["primary"]),
     ("¥105,000〜175,000", "獲得単価", "5件なら第1期平均と同水準", COLORS["primary"]),
     ("40%以上", "地元の集客構成比", "現在19%。成果構成比53%に近づける", COLORS["accent"]),
     ("40%以上", "DV360のサイト到達率", "現在15.9%。配信面の除外による改善", COLORS["accent"])]
for i, (v, l, n, ac) in enumerate(k):
    kpi(s, MX + i * (w4 + G), BODY_Y + 0.24, w4, 2.1, v, l, n,
        accent=ac, vcolor=ac, vsize=22 if len(v) > 8 else 30)
panel(s, MX, BODY_Y + 2.6, CW, 1.16, "運用体制", "週次でご報告し、8月中旬に配分を見直します",
      "8月30日の視察プレー受付終了に向け、月内で調整します。", accent=COLORS["primary"])
note_line(s, 6.1,
          "〈当社仮説〉KPIは第1期の実績水準を根拠とした目標値であり、達成を保証するものではありません。\n"
          "〈両論併記〉7月1日の3変更が同時発生していたため、施策単独の効果は8月の実測で検証します。")

# ══ 27  スケジュール・依頼事項 ═══════════════════
s = slide_new("5 提案", "スケジュールとご依頼事項", tag="C")
txt(s, MX, BODY_Y, 6.3, 0.3, "スケジュール", 14, COLORS["primary"], True)
rows = [
    ["期日", "内容"],
    ["7月下旬", "配信エリアの設定変更／DV360除外リスト作成／クラブ名検索の除外／コンバージョン計測の接続"],
    ["7月末", "クリエイティブを視察プレー訴求へ差し替え／LINEヤフー広告の入稿"],
    ["**8月1日", "**配信開始"],
    ["8月中旬", "中間レポート・配分の見直し"],
    ["**8月30日", "**視察プレー受付終了"],
    ["9月上旬", "8月実績のご報告・9月以降のご提案"],
]
table(s, MX, BODY_Y + 0.36, 6.3, rows, [1.4, 4.9],
      aligns=[PP_ALIGN.LEFT, PP_ALIGN.LEFT], hl=[3, 5], row_h=0.48, head_h=0.46, fs=10.5)
rx, rw = MX + 6.3 + G, CW - 6.3 - G
txt(s, rx, BODY_Y, rw, 0.3, "ご依頼事項", 14, COLORS["primary"], True)
rect(s, rx, BODY_Y + 0.36, rw, 3.62, COLORS["primary"])
txt(s, rx + 0.28, BODY_Y + 0.6, rw - 0.56, 3.2,
    "公式サイト経由のお問い合わせ実績を、\n月別・「知ったきっかけ」別でご共有ください。\n\n"
    + bullets([
        "公式サイトのフォームには「知ったきっかけ」が必須項目としてあり、「Googleの広告」「Facebook/Instagram等の広告」が選択肢に含まれています",
        "LPから公式サイトへ1,430件の遷移が発生しており、広告経由のお問い合わせが公式サイト側で受け付けられている可能性があります",
        "現時点の集計はLP経由のみです。実際の成果を過小評価している可能性があります",
    ]),
    SZ["body_s"], COLORS["bg"], line=1.55)
note_line(s, 6.28,
          "本レポートの出典：お問い合わせ通知メールの実データ／Google Analytics 4／Google Ads API／Meta Marketing API／"
          "DV360 Bid Manager API／LINEヤフー広告API／弊社基幹システム")

os.makedirs(OUT_DIR, exist_ok=True)
prs.save(OUT)
print(f"保存: {OUT}")
print(f"全{_st['page']}枚")
print("レイアウト順: " + " ".join(_st["layouts"]))
